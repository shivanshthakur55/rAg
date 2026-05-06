"""
RAG Engine — Universal document ingestion and conversational Q&A.

Supported formats:
  - PDF  (text-layer via PyPDF; falls back to Tesseract OCR for image/scanned PDFs)
  - DOCX / DOC  (via Docx2txt)
  - TXT          (plain text)
  - XLSX / XLS   (Excel — converted to markdown tables via pandas)
  - CSV          (comma-separated — converted to text via pandas)
  - PPTX         (PowerPoint — slide text extraction via python-pptx)

Design:
  - HuggingFaceEmbeddings (all-mpnet-base-v2)
  - Chroma persisted under Vector/pdf_documents/ (separate from invoices)
  - Each uploaded file gets its own Chroma collection (session_id)
  - A virtual "all" session searches across all persisted collections
  - ChatGroq (llama-3.3-70b-versatile) as the LLM
  - LCEL chain with message-history for follow-up question support
"""

import io
import uuid
from pathlib import Path
from typing import Optional

import pytesseract
from dotenv import load_dotenv

# Force Tesseract path for Windows
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

from langchain_chroma import Chroma
from langchain_community.document_loaders import PyPDFLoader, TextLoader, Docx2txtLoader
from langchain_core.documents import Document
from langchain_core.messages import AIMessage, HumanMessage
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_groq import ChatGroq
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from PIL import Image

load_dotenv()

# ── Constants ─────────────────────────────────────────────────────────────────
VECTOR_DIR = Path(__file__).parent.parent / "Vector" / "pdf_documents"
VECTOR_DIR.mkdir(parents=True, exist_ok=True)
CHUNK_SIZE = 800
CHUNK_OVERLAP = 200

# ── Singletons ────────────────────────────────────────────────────────────────
_embedding_model: Optional[HuggingFaceEmbeddings] = None
_llm: Optional[ChatGroq] = None


def get_embedding_model() -> HuggingFaceEmbeddings:
    global _embedding_model
    if _embedding_model is None:
        _embedding_model = HuggingFaceEmbeddings(
            model_name="sentence-transformers/all-mpnet-base-v2"
        )
    return _embedding_model


def get_llm() -> ChatGroq:
    global _llm
    if _llm is None:
        _llm = ChatGroq(
            model="llama-3.3-70b-versatile",
            temperature=0.7,
            max_tokens=4096,
        )
    return _llm


# ── In-memory session state ───────────────────────────────────────────────────
# { session_id: {"history": [...], "filename": str} }
_sessions: dict[str, dict] = {}


# ── Universal Document Loaders ────────────────────────────────────────────────

def _load_excel(file_path: str) -> list[Document]:
    """Convert every sheet of an Excel file into text using pandas."""
    import pandas as pd
    xl = pd.ExcelFile(file_path)
    docs = []
    for sheet in xl.sheet_names:
        df = xl.parse(sheet)
        # Drop fully-empty rows/cols for cleaner output
        df = df.dropna(how="all").dropna(axis=1, how="all")
        text = f"Sheet: {sheet}\n\n" + df.to_string(index=False)
        docs.append(Document(page_content=text, metadata={"sheet": sheet}))
    return docs


def _load_csv(file_path: str) -> list[Document]:
    """Load a CSV file as a human-readable text document."""
    import pandas as pd
    df = pd.read_csv(file_path)
    df = df.dropna(how="all").dropna(axis=1, how="all")
    text = df.to_string(index=False)
    return [Document(page_content=text)]


def _load_pptx(file_path: str) -> list[Document]:
    """Extract text from each slide of a PowerPoint file."""
    from pptx import Presentation
    prs = Presentation(file_path)
    docs = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
        if parts:
            docs.append(Document(
                page_content=f"Slide {i}:\n" + "\n".join(parts),
                metadata={"slide": i},
            ))
    return docs


def _ocr_pdf(file_path: str) -> list[Document]:
    """
    Render each PDF page to an image using PyMuPDF (no Poppler needed)
    then run Tesseract OCR on each page image.
    """
    import fitz  # PyMuPDF
    doc = fitz.open(file_path)
    result_docs = []
    try:
        for page_num in range(len(doc)):
            page = doc[page_num]
            mat = fitz.Matrix(2.0, 2.0)
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))
            text = pytesseract.image_to_string(img, lang="eng").strip()
            if text:
                result_docs.append(Document(
                    page_content=text,
                    metadata={"page": page_num, "source": file_path, "ocr": True},
                ))
    finally:
        doc.close()
    return result_docs


def _ocr_image(file_path: str) -> list[Document]:
    """Run Tesseract OCR directly on an image file (JPG, PNG, TIFF, BMP, WEBP)."""
    img = Image.open(file_path)
    text = pytesseract.image_to_string(img, lang="eng").strip()
    if not text:
        return []
    return [Document(page_content=text, metadata={"source": file_path, "ocr": True})]


def load_document(file_path: str) -> list[Document]:
    """
    Universal document loader — dispatches to the appropriate handler
    based on file extension. Returns a list of LangChain Documents.
    """
    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        # Try PyPDF first (fast, no OCR needed for digital PDFs)
        try:
            docs = PyPDFLoader(file_path).load()
            # Check if meaningful text was extracted
            total_text = "".join(d.page_content for d in docs).strip()
            if len(total_text) < 50:
                # Very little text → likely scanned/image PDF → fall back to OCR
                raise ValueError("Insufficient text in PyPDF output — trying OCR")
            return docs
        except Exception:
            # OCR fallback
            return _ocr_pdf(file_path)

    elif ext == ".txt":
        return TextLoader(file_path, encoding="utf-8").load()

    elif ext in (".docx", ".doc"):
        return Docx2txtLoader(file_path).load()

    elif ext in (".xlsx", ".xls"):
        return _load_excel(file_path)

    elif ext == ".csv":
        return _load_csv(file_path)

    elif ext in (".pptx", ".ppt"):
        return _load_pptx(file_path)

    elif ext in (".jpg", ".jpeg", ".png", ".tiff", ".bmp", ".webp"):
        # Direct image OCR
        return _ocr_image(file_path)

    else:
        raise ValueError(f"Unsupported file type: {ext}")


# ── Document ingestion ────────────────────────────────────────────────────────

def ingest_document(file_path: str, filename: str) -> str:
    """Load → chunk → embed → persist in Chroma. Returns session_id."""
    session_id = str(uuid.uuid4())
    docs = load_document(file_path)

    for doc in docs:
        doc.metadata["uploaded_filename"] = filename
        doc.metadata["session_id"] = session_id

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP
    )
    chunks = splitter.split_documents(docs)

    Chroma.from_documents(
        documents=chunks,
        embedding=get_embedding_model(),
        collection_name=session_id,
        persist_directory=str(VECTOR_DIR),
    )

    _sessions[session_id] = {"history": [], "filename": filename}
    return session_id


# ── Vector store helpers ──────────────────────────────────────────────────────

def _get_retriever(session_id: str, k: int = 8):
    vs = Chroma(
        collection_name=session_id,
        embedding_function=get_embedding_model(),
        persist_directory=str(VECTOR_DIR),
    )
    return vs.as_retriever(search_kwargs={"k": k})


def _get_all_session_ids() -> list[str]:
    """Return all session_ids that have a persisted Chroma collection."""
    import chromadb
    client = chromadb.PersistentClient(path=str(VECTOR_DIR))
    # Filter out any non-UUID collection names (e.g. stray collections)
    collections = [c.name for c in client.list_collections()]
    return collections


# ── Prompt templates ──────────────────────────────────────────────────────────

_contextualize_prompt = ChatPromptTemplate.from_messages([
    ("system",
     "Given the chat history and the latest user question, "
     "reformulate the question to be standalone so it can be understood "
     "without the history. Do NOT answer. Just reformulate if needed, else return as is."),
    MessagesPlaceholder("chat_history"),
    ("human", "{input}"),
])

_qa_prompt = ChatPromptTemplate.from_messages([
    ("system",
     "Expert research assistant. Answer ONLY using provided context. No outside info.\n\n"
     "- Use bullet points for clarity.\n"
     "- Be concise: provide short, direct answers. Ask 'Do you need more information?' and stop. Provide thorough details ONLY if the user says yes.\n"
     "- Quote context to support answers. Address all query aspects.\n"
     "- If the answer is missing, state what is available in the context instead.\n"
     "- If the context comes from multiple files, mention which file each piece of information is from.\n\n"
     "Context:\n{context}"),
    MessagesPlaceholder("chat_history"),
    ("human", "{input}"),
])


def _format_docs(docs: list) -> str:
    parts = []
    for d in docs:
        fname = d.metadata.get("uploaded_filename", d.metadata.get("source", ""))
        page = d.metadata.get("page", None)
        header = f"[Source: {fname}" + (f", page {page + 1}" if page is not None else "") + "]"
        parts.append(f"{header}\n{d.page_content}")
    return "\n\n".join(parts)


# ── RAG chain ─────────────────────────────────────────────────────────────────

def _run_rag(session_id: str, question: str) -> tuple[str, list]:
    """
    Single-session RAG: retrieve from one Chroma collection and generate answer.
    Returns (answer_str, source_docs).
    """
    llm = get_llm()
    retriever = _get_retriever(session_id)
    history = _sessions[session_id]["history"]

    # Rephrase follow-up if history exists
    if history:
        rephrased = (
            _contextualize_prompt | llm | StrOutputParser()
        ).invoke({"input": question, "chat_history": history})
    else:
        rephrased = question

    docs = retriever.invoke(rephrased)
    context = _format_docs(docs)

    answer = (
        _qa_prompt | llm | StrOutputParser()
    ).invoke({"input": question, "context": context, "chat_history": history})

    return answer, docs


def _run_rag_all(question: str, history: list) -> tuple[str, list]:
    """
    Cross-document RAG: retrieve from ALL persisted Chroma collections,
    merge chunks, deduplicate, and generate one unified answer.
    """
    llm = get_llm()

    # Rephrase if there's history
    if history:
        rephrased = (
            _contextualize_prompt | llm | StrOutputParser()
        ).invoke({"input": question, "chat_history": history})
    else:
        rephrased = question

    # Gather chunks from every collection
    all_docs: list = []
    seen_content: set[str] = set()
    collection_ids = _get_all_session_ids()

    for sid in collection_ids:
        try:
            vs = Chroma(
                collection_name=sid,
                embedding_function=get_embedding_model(),
                persist_directory=str(VECTOR_DIR),
            )
            results = vs.similarity_search(rephrased, k=4)
            for doc in results:
                key = doc.page_content[:120]  # rough dedup key
                if key not in seen_content:
                    seen_content.add(key)
                    all_docs.append(doc)
        except Exception:
            continue  # skip broken collections silently

    if not all_docs:
        return "No documents have been uploaded yet. Please upload a file first.", []

    context = _format_docs(all_docs)
    answer = (
        _qa_prompt | llm | StrOutputParser()
    ).invoke({"input": question, "context": context, "chat_history": history})

    return answer, all_docs


# ── Public API ────────────────────────────────────────────────────────────────

def chat(session_id: str, question: str) -> dict:
    """Chat against a specific document session."""
    if session_id not in _sessions:
        # Try to restore from persisted Chroma collection
        try:
            Chroma(
                collection_name=session_id,
                embedding_function=get_embedding_model(),
                persist_directory=str(VECTOR_DIR),
            ).get()
            _sessions[session_id] = {"history": [], "filename": "unknown"}
        except Exception:
            raise ValueError(f"Unknown session: {session_id}")

    answer, docs = _run_rag(session_id, question)

    _sessions[session_id]["history"].append(HumanMessage(content=question))
    _sessions[session_id]["history"].append(AIMessage(content=answer))

    sources = list({
        doc.metadata.get("uploaded_filename", doc.metadata.get("source", ""))
        + (f" (page {doc.metadata['page'] + 1})" if "page" in doc.metadata else "")
        for doc in docs
    })

    return {"answer": answer, "sources": sources}


def chat_all(question: str) -> dict:
    """
    Search across ALL uploaded documents (no session selection needed).
    Uses a shared 'all' history key in _sessions.
    """
    ALL_KEY = "__all__"
    if ALL_KEY not in _sessions:
        _sessions[ALL_KEY] = {"history": [], "filename": "All Documents"}

    history = _sessions[ALL_KEY]["history"]
    answer, docs = _run_rag_all(question, history)

    _sessions[ALL_KEY]["history"].append(HumanMessage(content=question))
    _sessions[ALL_KEY]["history"].append(AIMessage(content=answer))

    sources = list({
        doc.metadata.get("uploaded_filename", doc.metadata.get("source", ""))
        + (f" (page {doc.metadata['page'] + 1})" if "page" in doc.metadata else "")
        for doc in docs
    })

    return {"answer": answer, "sources": sources}


def clear_all_history() -> None:
    """Clear the cross-document search history."""
    ALL_KEY = "__all__"
    if ALL_KEY in _sessions:
        _sessions[ALL_KEY]["history"] = []


def list_sessions() -> list[dict]:
    return [
        {"session_id": sid, "filename": data["filename"]}
        for sid, data in _sessions.items()
        if sid != "__all__"
    ]


def clear_session_history(session_id: str):
    if session_id in _sessions:
        _sessions[session_id]["history"] = []
